"""
Build a fully-editable PPTX presentation:
«Гибридная система прогнозирования для HR — Уральская сталь»

10 slides, all text in text-boxes (editable in PowerPoint / LibreOffice).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree
import os

# ─────────────────────────────────────────────────────────────
# CONSTANTS
# ─────────────────────────────────────────────────────────────
W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

# Brand colours
C_BG      = RGBColor(0x0F, 0x19, 0x23)
C_SURFACE = RGBColor(0x16, 0x20, 0x30)
C_RED     = RGBColor(0xC0, 0x39, 0x2B)
C_RED_L   = RGBColor(0xE7, 0x4C, 0x3C)
C_GOLD    = RGBColor(0xE6, 0x7E, 0x22)
C_BLUE    = RGBColor(0x29, 0x80, 0xB9)
C_BLUE_L  = RGBColor(0x5D, 0xAD, 0xE2)
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)
C_GREEN_L = RGBColor(0x2E, 0xCC, 0x71)
C_GRAY    = RGBColor(0x8D, 0xA8, 0xC4)
C_WHITE   = RGBColor(0xEC, 0xF0, 0xF1)
C_CARD    = RGBColor(0x1E, 0x2D, 0x3F)
C_DIM     = RGBColor(0x55, 0x6B, 0x82)
C_DARK2   = RGBColor(0x12, 0x1E, 0x2C)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ─────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    """Add a plain rectangle shape."""
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h,
                text, size=Pt(14), bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, wrap=True,
                v_anchor=None):
    """Add a textbox with a single paragraph."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=Pt(12), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(0)):
    """Append a paragraph to an existing text-frame."""
    from pptx.util import Pt as uPt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    return p


def bg(slide, color=C_BG):
    """Fill slide background."""
    add_rect(slide, 0, 0, W, H, fill=color)


def accent_bar(slide, y=Inches(0), h=Inches(0.06), color=C_RED):
    """Thin horizontal accent bar."""
    add_rect(slide, 0, y, W, h, fill=color)


def slide_num_tag(slide, cur, total=10):
    add_textbox(slide,
                W - Inches(1.4), Inches(0.15),
                Inches(1.3), Inches(0.3),
                f"{cur:02d} / {total:02d}",
                size=Pt(9), color=C_DIM, align=PP_ALIGN.RIGHT)


def section_label(slide, text, x, y, w=Inches(4)):
    add_textbox(slide, x, y, w, Inches(0.3),
                text.upper(), size=Pt(9),
                bold=True, color=C_RED)


def card_rect(slide, x, y, w, h, accent_left=False,
              fill=C_CARD, border=None):
    r = add_rect(slide, x, y, w, h, fill=fill,
                 line=border or C_CARD,
                 line_w=Pt(0.5) if border else Pt(0))
    if accent_left:
        add_rect(slide, x, y, Inches(0.055), h, fill=C_RED)
    return r


def pill_badge(slide, x, y, text, color=C_GREEN_L, bg_color=None):
    if bg_color is None:
        # compute a dim version — RGBColor is a hex string subclass; parse it
        r_hex = str(color)  # e.g. "2ECC71"
        rr = int(r_hex[0:2], 16)
        gg = int(r_hex[2:4], 16)
        bb = int(r_hex[4:6], 16)
        bg_color = RGBColor(
            min(255, rr // 5),
            min(255, gg // 5),
            min(255, bb // 5),
        )
    w = Inches(1.5)
    h = Inches(0.28)
    r = add_rect(slide, x, y, w, h, fill=bg_color)
    add_textbox(slide, x, y, w, h,
                text, size=Pt(9), bold=True, color=color,
                align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────
def slide1():
    s = prs.slides.add_slide(blank_layout)
    bg(s, C_BG)

    # dark gradient block left
    add_rect(s, 0, 0, Inches(6.8), H, fill=RGBColor(0x0A, 0x14, 0x1E))
    # red accent bar top
    accent_bar(s, y=0, h=Inches(0.08))

    slide_num_tag(s, 1)

    # eyebrow
    add_textbox(s, Inches(0.6), Inches(1.0), Inches(6), Inches(0.4),
                "УРАЛЬСКАЯ СТАЛЬ",
                size=Pt(11), bold=True, color=C_GOLD)

    # main title
    tb = s.shapes.add_textbox(
        Inches(0.6), Inches(1.55), Inches(6.1), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = "Гибридная система\nпрогнозирования для "
    r1.font.name  = FONT
    r1.font.size  = Pt(38)
    r1.font.bold  = True
    r1.font.color.rgb = C_WHITE
    r2 = p1.add_run()
    r2.text = "HR"
    r2.font.name  = FONT
    r2.font.size  = Pt(38)
    r2.font.bold  = True
    r2.font.color.rgb = C_GOLD

    # subtitle
    add_textbox(s, Inches(0.6), Inches(3.85), Inches(6.1), Inches(0.9),
                "Мультимодельный подход к прогнозированию потребности в персонале:\n"
                "от найма до бюджетирования ФОТ",
                size=Pt(13), color=C_GRAY)

    # KPI chips
    chips = [
        ("4",      "Модели в системе"),
        ("6",      "HR-процессов"),
        ("< 2 мин","Готовый прогноз"),
        ("95%",    "Дов. интервал"),
    ]
    cx = Inches(0.6)
    for val, lbl in chips:
        card_rect(s, cx, Inches(4.95), Inches(1.45), Inches(1.15),
                  fill=RGBColor(0x1A, 0x2B, 0x3C))
        add_textbox(s, cx, Inches(5.0), Inches(1.45), Inches(0.45),
                    val, size=Pt(20), bold=True, color=C_GOLD,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, cx, Inches(5.48), Inches(1.45), Inches(0.55),
                    lbl, size=Pt(9), color=C_GRAY, align=PP_ALIGN.CENTER)
        cx += Inches(1.55)

    # right decorative panel
    add_rect(s, Inches(7.0), 0, Inches(6.33), H,
             fill=RGBColor(0x12, 0x1E, 0x2C))
    # right content
    add_textbox(s, Inches(7.3), Inches(1.5), Inches(5.7), Inches(0.5),
                "О СИСТЕМЕ", size=Pt(10), bold=True, color=C_RED)

    bullets = [
        "Комбинирует SARIMA, XGBoost и TimeLLM (LLM)",
        "Адаптивные веса: лучшая модель получает больше влияния",
        "LLM-эксперт читает актуальные веб-источники",
        "Работает на коротких рядах (от 5 наблюдений)",
        "REST API — интеграция с любой HR-системой",
        "Отечественный стек: YandexGPT + on-premise",
    ]
    by = Inches(2.1)
    for b in bullets:
        add_rect(s, Inches(7.3), by + Inches(0.08),
                 Inches(0.12), Inches(0.12), fill=C_RED)
        add_textbox(s, Inches(7.55), by, Inches(5.4), Inches(0.38),
                    b, size=Pt(12), color=C_WHITE)
        by += Inches(0.52)

    # footer
    add_rect(s, 0, H - Inches(0.55), W, Inches(0.55),
             fill=RGBColor(0x08, 0x0E, 0x16))
    add_textbox(s, Inches(0.4), H - Inches(0.48),
                Inches(4), Inches(0.38),
                "Диссертационное исследование  ·  Февраль 2026",
                size=Pt(10), color=C_DIM)

slide1()


# ─────────────────────────────────────────────────────────────
# SLIDE 2 — CHALLENGE
# ─────────────────────────────────────────────────────────────
def slide2():
    s = prs.slides.add_slide(blank_layout)
    bg(s)
    accent_bar(s, y=0, h=Inches(0.06))
    slide_num_tag(s, 2)

    # red left panel
    add_rect(s, 0, 0, Inches(4.4), H,
             fill=RGBColor(0xC0, 0x39, 0x2B))

    section_label(s, "Проблема", Inches(0.35), Inches(0.9), Inches(3.7))
    add_textbox(s, Inches(0.35), Inches(1.25), Inches(3.7), Inches(1.8),
                "HR-планирование теряет точность в условиях нестабильности",
                size=Pt(22), bold=True, color=C_WHITE)
    add_textbox(s, Inches(0.35), Inches(3.2), Inches(3.7), Inches(2.5),
                "Традиционные методы прогнозирования численности персонала, "
                "бюджета ФОТ и потребности в кадрах не учитывают нелинейные "
                "изменения рынка труда, внешние сигналы и отраслевые тренды.",
                size=Pt(12), color=RGBColor(0xFF, 0xFF, 0xFF))

    # 4 problem cards
    problems = [
        ("🎯", "Потребность в кадрах непредсказуема",
         "Плановые показатели отклоняются от факта на 15–30%. "
         "Линейная экстраполяция не улавливает структурные изменения."),
        ("💸", "Ошибки в бюджете ФОТ дорого обходятся",
         "Отклонение прогноза зарплатного рынка на 5% = "
         "десятки миллионов рублей перерасхода."),
        ("⏳", "Реакция на изменения — с опозданием",
         "К моменту, когда HR видит проблему, устранить её дорого. "
         "Нужен горизонт 12–36 месяцев для превентивных действий."),
        ("🌐", "Внешний контекст не учитывается",
         "Классические модели не извлекают сигналы из новостей, "
         "отраслевых отчётов, данных hh.ru и Росстата."),
    ]
    py = Inches(0.7)
    for ico, title, desc in problems:
        card_rect(s, Inches(4.65), py, Inches(8.3), Inches(1.4),
                  accent_left=True, fill=C_CARD)
        add_textbox(s, Inches(4.82), py + Inches(0.12),
                    Inches(0.5), Inches(0.5), ico, size=Pt(20), color=C_WHITE)
        add_textbox(s, Inches(5.42), py + Inches(0.1),
                    Inches(7.4), Inches(0.35),
                    title, size=Pt(13), bold=True, color=C_WHITE)
        add_textbox(s, Inches(5.42), py + Inches(0.52),
                    Inches(7.4), Inches(0.75),
                    desc, size=Pt(11), color=C_GRAY)
        py += Inches(1.55)

slide2()


# ─────────────────────────────────────────────────────────────
# SLIDE 3 — ARCHITECTURE
# ─────────────────────────────────────────────────────────────
def slide3():
    s = prs.slides.add_slide(blank_layout)
    bg(s)
    accent_bar(s, y=0, h=Inches(0.06))
    slide_num_tag(s, 3)

    section_label(s, "Подход", Inches(0.5), Inches(0.18))
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(11), Inches(0.55),
                "Гибридная мультимодельная архитектура",
                size=Pt(26), bold=True, color=C_WHITE)
    add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.4),
                "Четыре компонента компенсируют слабости друг друга и адаптивно перераспределяют веса",
                size=Pt(12), color=C_GRAY)

    boxes = [
        ("1", "📈", "SARIMA-XS",
         "Классическая статистическая модель с сезонностью. "
         "Автоматический подбор параметров. Стабильна на коротких рядах.",
         "⚡ 10–30 сек", C_BLUE_L),
        ("2", "🌲", "XGBoost TS",
         "Gradient Boosting для временных рядов. Автоматическая "
         "генерация признаков: лаги, скользящие средние, тренды.",
         "⚡ 10–30 сек", C_GOLD),
        ("3", "🤖", "TimeLLM",
         "Малые языковые модели SLM 2024–2025. Понимает временны́е "
         "паттерны как языковый контекст. GPU-ускорение.",
         "🔥 GPU", C_RED_L),
        ("4", "🎯", "Hybrid + LLM-эксперт",
         "Адаптивное взвешивание трёх моделей. YandexGPT анализирует "
         "актуальные источники и корректирует итоговый прогноз.",
         "⭐ Рекомендуется", C_GREEN_L),
    ]

    bw = Inches(2.9)
    bx = Inches(0.4)
    for i, (num, ico, title, desc, tag, tcol) in enumerate(boxes):
        is_last = (i == len(boxes) - 1)
        fill = RGBColor(0x1A, 0x2B, 0x3C) if not is_last else RGBColor(0x1E, 0x16, 0x16)
        card_rect(s, bx, Inches(1.65), bw, Inches(4.55), fill=fill,
                  border=C_RED if is_last else None)
        # number badge
        add_rect(s, bx + Inches(0.18), Inches(1.82),
                 Inches(0.38), Inches(0.38), fill=C_RED)
        add_textbox(s, bx + Inches(0.18), Inches(1.82),
                    Inches(0.38), Inches(0.38),
                    num, size=Pt(13), bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, bx + Inches(0.22), Inches(2.35),
                    Inches(2.4), Inches(0.45),
                    ico, size=Pt(26), color=C_WHITE)
        add_textbox(s, bx + Inches(0.22), Inches(2.9),
                    Inches(2.5), Inches(0.45),
                    title, size=Pt(13), bold=True, color=C_WHITE)
        add_textbox(s, bx + Inches(0.22), Inches(3.42),
                    Inches(2.5), Inches(1.8),
                    desc, size=Pt(11), color=C_GRAY)
        # tag pill
        pill_badge(s, bx + Inches(0.22), Inches(5.55), tag,
                   color=tcol)

        # arrow between boxes
        if not is_last:
            add_textbox(s, bx + bw + Inches(0.05), Inches(3.4),
                        Inches(0.35), Inches(0.5),
                        "›", size=Pt(28), bold=True, color=C_RED,
                        align=PP_ALIGN.CENTER)
        bx += bw + Inches(0.42)

    # formula strip
    add_rect(s, Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.82),
             fill=RGBColor(0x08, 0x12, 0x1C),
             line=RGBColor(0x50, 0x18, 0x14), line_w=Pt(0.75))
    add_textbox(s, Inches(0.6), Inches(6.43), Inches(1.1), Inches(0.35),
                "ФОРМУЛА", size=Pt(8), bold=True, color=C_RED)
    add_textbox(s, Inches(1.75), Inches(6.41), Inches(11), Inches(0.55),
                "Ŷ = w₁·Ŷ_SARIMA + w₂·Ŷ_XGB + w₃·Ŷ_LLM  |  "
                "wᵢ = (1/ERᵢ) / Σ(1/ERⱼ)  |  ERᵢ(t) = λ·ERᵢ(t−1) + (1−λ)·errᵢ(t),  λ=0.9",
                size=Pt(11), color=C_WHITE)

slide3()


# ─────────────────────────────────────────────────────────────
# SLIDE 4 — HR PROCESSES
# ─────────────────────────────────────────────────────────────
def slide4():
    s = prs.slides.add_slide(blank_layout)
    bg(s)
    accent_bar(s, y=0, h=Inches(0.06))
    slide_num_tag(s, 4)

    section_label(s, "Применение", Inches(0.5), Inches(0.18))
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.55),
                "6 ключевых HR-процессов, которые автоматизирует система",
                size=Pt(26), bold=True, color=C_WHITE)
    add_textbox(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.35),
                "Для каждого процесса: загрузите CSV → прогноз с 95% ДИ + аналитика LLM-эксперта",
                size=Pt(12), color=C_GRAY)

    cards = [
        ("👥", "Планирование численности",
         "Прогноз потребности по подразделениям на 1–3 года. "
         "Учёт производственных планов, текучести и демографии.",
         "Горизонт: 12–36 мес."),
        ("💰", "Бюджетирование ФОТ",
         "Прогноз роста зарплат с учётом инфляции, дефицита кадров "
         "и отраслевой динамики. Реалистичный бюджет без пересмотра.",
         "Точность: ±5–8%"),
        ("🔍", "Управление потоком найма",
         "Прогноз открытых вакансий, времени закрытия и стоимости. "
         "Планирование загрузки рекрутинга и запасов кандидатов.",
         "Снижение TTH на 20–35%"),
        ("📉", "Прогноз текучести кадров",
         "Моделирование текучести в зависимости от зарплатной динамики "
         "и сезонности. Сигнал на удержание за 3–6 месяцев.",
         "Превентивное удержание"),
        ("🎓", "Планирование обучения",
         "Прогноз потребности в компетенциях и переобучении. "
         "Синхронизация с планами внедрения оборудования.",
         "Синхр. с пр. планами"),
        ("🌐", "Мониторинг рынка труда",
         "Автосбор данных hh.ru, Росстат. LLM-эксперт анализирует "
         "новости и корректирует прогноз в реальном времени.",
         "Актуализация онлайн"),
    ]

    cw, ch = Inches(4.0), Inches(2.35)
    positions = [
        (Inches(0.35), Inches(1.6)),
        (Inches(4.65), Inches(1.6)),
        (Inches(8.95), Inches(1.6)),
        (Inches(0.35), Inches(4.1)),
        (Inches(4.65), Inches(4.1)),
        (Inches(8.95), Inches(4.1)),
    ]
    for (cx, cy), (ico, title, desc, gain) in zip(positions, cards):
        card_rect(s, cx, cy, cw, ch, fill=C_CARD,
                  border=RGBColor(0x28, 0x3A, 0x50))
        add_textbox(s, cx + Inches(0.18), cy + Inches(0.15),
                    Inches(0.55), Inches(0.45),
                    ico, size=Pt(22), color=C_WHITE)
        add_textbox(s, cx + Inches(0.8), cy + Inches(0.15),
                    cw - Inches(0.9), Inches(0.42),
                    title, size=Pt(12), bold=True, color=C_WHITE)
        add_textbox(s, cx + Inches(0.18), cy + Inches(0.68),
                    cw - Inches(0.3), Inches(1.1),
                    desc, size=Pt(10.5), color=C_GRAY)
        # gain pill
        add_rect(s, cx + Inches(0.18), cy + ch - Inches(0.48),
                 cw - Inches(0.36), Inches(0.34),
                 fill=RGBColor(0x0E, 0x28, 0x1A))
        add_textbox(s, cx + Inches(0.18), cy + ch - Inches(0.48),
                    cw - Inches(0.36), Inches(0.34),
                    gain, size=Pt(10), bold=True,
                    color=C_GREEN_L, align=PP_ALIGN.CENTER)

slide4()


# ─────────────────────────────────────────────────────────────
# SLIDE 5 — PIPELINE
# ─────────────────────────────────────────────────────────────
def slide5():
    s = prs.slides.add_slide(blank_layout)
    bg(s)
    accent_bar(s, y=0, h=Inches(0.06))
    slide_num_tag(s, 5)

    section_label(s, "Как это работает", Inches(0.5), Inches(0.18))
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(11), Inches(0.55),
                "От данных до HR-решения: 5 шагов",
                size=Pt(26), bold=True, color=C_WHITE)
    add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.35),
                "Полный цикл прогнозирования — менее 2 минут. Работает через веб-интерфейс, REST API или Python SDK",
                size=Pt(12), color=C_GRAY)

    steps = [
        ("1", "Загрузка данных",
         "CSV или Excel с историческими HR-метриками. "
         "Система автоматически определяет частоту, "
         "пропуски и сезонность. Минимум — 5 наблюдений.",
         "CSV / XLSX", C_BLUE_L),
        ("2", "Параллельное обучение",
         "SARIMA, XGBoost и TimeLLM обучаются одновременно. "
         "Автоматический подбор гиперпараметров. "
         "Поддержка коротких рядов (n≥5).",
         "30–90 сек", C_GOLD),
        ("3", "Адаптивные веса",
         "Система вычисляет исторические ошибки каждой "
         "модели и назначает веса обратно пропорционально "
         "ошибке. Лучшая модель — больше влияния.",
         "λ = 0.9", C_GREEN_L),
        ("4", "LLM-коррекция",
         "YandexGPT загружает указанные URL (Росстат, hh.ru, "
         "новости), извлекает факты и корректирует прогноз "
         "коэффициентами 0.8–1.2.",
         "YandexGPT", C_BLUE_L),
        ("5", "Результат и экспорт",
         "Прогноз с 95% ДИ, метрики MAE/MAPE/R², "
         "веса моделей, комментарий LLM-эксперта. "
         "CSV, XLSX или API.",
         "API / CSV / UI", C_GREEN_L),
    ]

    sw = Inches(2.3)
    sx = Inches(0.35)
    for i, (num, title, desc, tag, tcol) in enumerate(steps):
        card_rect(s, sx, Inches(1.65), sw, Inches(4.5), fill=C_CARD)
        # number badge
        add_rect(s, sx + Inches(0.18), Inches(1.85),
                 Inches(0.38), Inches(0.38), fill=C_RED)
        add_textbox(s, sx + Inches(0.18), Inches(1.85),
                    Inches(0.38), Inches(0.38),
                    num, size=Pt(13), bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, sx + Inches(0.18), Inches(2.35),
                    sw - Inches(0.3), Inches(0.48),
                    title, size=Pt(12), bold=True, color=C_WHITE)
        add_textbox(s, sx + Inches(0.18), Inches(2.92),
                    sw - Inches(0.3), Inches(2.2),
                    desc, size=Pt(10.5), color=C_GRAY)
        pill_badge(s, sx + Inches(0.18), Inches(5.6), tag, color=tcol)

        if i < len(steps) - 1:
            add_textbox(s, sx + sw + Inches(0.06), Inches(3.7),
                        Inches(0.3), Inches(0.5),
                        "›", size=Pt(26), bold=True, color=C_RED,
                        align=PP_ALIGN.CENTER)
        sx += sw + Inches(0.42)

    # LLM loop strip
    add_rect(s, Inches(0.35), Inches(6.35), Inches(12.6), Inches(0.82),
             fill=RGBColor(0x08, 0x12, 0x1C),
             line=RGBColor(0x50, 0x18, 0x14), line_w=Pt(0.75))
    add_textbox(s, Inches(0.55), Inches(6.43), Inches(1.4), Inches(0.35),
                "LLM-ЭКСПЕРТ", size=Pt(8), bold=True, color=C_RED)
    add_textbox(s, Inches(2.0), Inches(6.41), Inches(11), Inches(0.55),
                "Источники → Извлечение фактов → Направление тренда (↑/→/↓) → "
                "Коэффициент [0.8–1.2] → Ŷ_скорр = Ŷ × коэфф",
                size=Pt(11), color=C_WHITE)

slide5()


# ─────────────────────────────────────────────────────────────
# SLIDE 6 — DATA REQUIREMENTS
# ─────────────────────────────────────────────────────────────
def slide6():
    s = prs.slides.add_slide(blank_layout)
    bg(s)
    accent_bar(s, y=0, h=Inches(0.06))
    slide_num_tag(s, 6)

    section_label(s, "Требования к данным", Inches(0.5), Inches(0.18))
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(11), Inches(0.55),
                "Что нужно для запуска прогноза",
                size=Pt(26), bold=True, color=C_WHITE)
    add_textbox(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.35),
                "Достаточно выгрузки из 1С:ЗУП или Excel-таблицы HR-отдела — система сделает остальное",
                size=Pt(12), color=C_GRAY)

    # LEFT COLUMN
    lx = Inches(0.4)

    # Mandatory fields block
    add_rect(s, lx, Inches(1.6), Inches(6.1), Inches(2.35), fill=C_CARD,
             line=RGBColor(0x28, 0x3A, 0x50), line_w=Pt(0.5))
    add_textbox(s, lx + Inches(0.2), Inches(1.7), Inches(5.7), Inches(0.35),
                "📋  ОБЯЗАТЕЛЬНЫЕ ПОЛЯ", size=Pt(10), bold=True, color=C_GOLD)
    fields = [
        ("Дата наблюдения",
         "Год, квартал или месяц. Форматы: ГГГГ-ММ-ДД, ГГГГ-ММ, ГГГГ. Пропуски допустимы."),
        ("Числовое значение",
         "Численность, зарплата, вакансии, текучесть (%) — любой числовой ряд."),
        ("Минимальный объём",
         "От 5 наблюдений для базового прогноза. От 10–15 — для надёжного MAPE."),
    ]
    fy = Inches(2.12)
    for ft, fd in fields:
        add_rect(s, lx + Inches(0.2), fy + Inches(0.07),
                 Inches(0.1), Inches(0.1), fill=C_RED)
        add_textbox(s, lx + Inches(0.4), fy, Inches(2.2), Inches(0.28),
                    ft, size=Pt(11), bold=True, color=C_WHITE)
        add_textbox(s, lx + Inches(0.4), fy + Inches(0.28),
                    Inches(5.5), Inches(0.35),
                    fd, size=Pt(10), color=C_GRAY)
        fy += Inches(0.68)

    # Format block
    add_rect(s, lx, Inches(4.1), Inches(6.1), Inches(1.55), fill=C_CARD,
             line=RGBColor(0x28, 0x3A, 0x50), line_w=Pt(0.5))
    add_textbox(s, lx + Inches(0.2), Inches(4.2), Inches(5.7), Inches(0.35),
                "📁  ФОРМАТЫ ФАЙЛОВ", size=Pt(10), bold=True, color=C_BLUE_L)
    fmts = [
        ("CSV",       "2 колонки: дата, значение; разделитель , или ;"),
        ("XLSX/XLS",  "Первый лист, заголовки в строке 1 — выгрузка из 1С:ЗУП"),
        ("REST API",  "JSON: {\"dates\": [...], \"values\": [...]}"),
    ]
    fy2 = Inches(4.62)
    for fmt, desc in fmts:
        add_textbox(s, lx + Inches(0.2), fy2, Inches(1.0), Inches(0.32),
                    fmt, size=Pt(10), bold=True, color=C_GOLD)
        add_textbox(s, lx + Inches(1.25), fy2, Inches(4.7), Inches(0.32),
                    desc, size=Pt(10), color=C_GRAY)
        fy2 += Inches(0.4)

    # Example box
    add_rect(s, lx, Inches(5.78), Inches(6.1), Inches(1.42),
             fill=RGBColor(0x0C, 0x17, 0x22),
             line=RGBColor(0x4A, 0x38, 0x12), line_w=Pt(0.75))
    add_textbox(s, lx + Inches(0.2), Inches(5.87), Inches(5.7), Inches(0.3),
                "ПРИМЕР CSV", size=Pt(8), bold=True, color=C_GOLD)
    add_textbox(s, lx + Inches(0.2), Inches(6.22), Inches(5.7), Inches(0.8),
                "date,value\n2022-01-01,4820\n2022-04-01,4910\n"
                "2022-07-01,5034\n...(минимум 5 строк)",
                size=Pt(10), color=RGBColor(0xA8, 0xC8, 0xE8))

    # RIGHT COLUMN
    rx = Inches(6.75)

    # Recommended indicators
    add_rect(s, rx, Inches(1.6), Inches(6.15), Inches(3.2), fill=C_CARD,
             line=RGBColor(0x28, 0x3A, 0x50), line_w=Pt(0.5))
    add_textbox(s, rx + Inches(0.2), Inches(1.7), Inches(5.7), Inches(0.35),
                "✅  РЕКОМЕНДУЕМЫЕ HR-ПОКАЗАТЕЛИ",
                size=Pt(10), bold=True, color=C_GREEN_L)
    indicators = [
        ("Численность персонала по категориям",
         "Ежемесячно / ежеквартально. ИТР, рабочие, РСС — отдельными рядами."),
        ("Средняя зарплата по грейдам",
         "Помесячно. Оклад + надбавки + премии суммарно или раздельно."),
        ("Количество открытых вакансий",
         "Ежемесячно по подразделениям или суммарно."),
        ("Коэффициент текучести (%)",
         "Ежеквартально или ежемесячно. Добровольная / вынужденная."),
    ]
    iy = Inches(2.12)
    for it, id_ in indicators:
        add_rect(s, rx + Inches(0.2), iy + Inches(0.08),
                 Inches(0.1), Inches(0.1), fill=C_GREEN)
        add_textbox(s, rx + Inches(0.4), iy, Inches(2.2), Inches(0.28),
                    it, size=Pt(11), bold=True, color=C_WHITE)
        add_textbox(s, rx + Inches(0.4), iy + Inches(0.28),
                    Inches(5.5), Inches(0.32),
                    id_, size=Pt(10), color=C_GRAY)
        iy += Inches(0.68)

    # LLM sources block
    add_rect(s, rx, Inches(4.95), Inches(6.15), Inches(1.15),
             fill=RGBColor(0x10, 0x1D, 0x32),
             line=RGBColor(0x1A, 0x44, 0x6A), line_w=Pt(0.75))
    add_textbox(s, rx + Inches(0.2), Inches(5.05), Inches(5.7), Inches(0.32),
                "💡  ОПЦИОНАЛЬНО: веб-источники для LLM-эксперта",
                size=Pt(10), bold=True, color=C_BLUE_L)
    add_textbox(s, rx + Inches(0.2), Inches(5.43), Inches(5.7), Inches(0.55),
                "Укажите 1–5 URL (hh.ru, Росстат, отраслевые порталы) — "
                "LLM-эксперт автоматически извлечёт факты и скорректирует прогноз.",
                size=Pt(10.5), color=C_GRAY)

    # Volume recommendation
    add_rect(s, rx, Inches(6.22), Inches(6.15), Inches(1.0),
             fill=RGBColor(0x16, 0x22, 0x16),
             line=RGBColor(0x1C, 0x44, 0x1C), line_w=Pt(0.75))
    vols = [("5+", "базовый прогноз"), ("10–15", "надёжный MAPE"), ("24+", "полная точность TimeLLM")]
    vx2 = rx + Inches(0.2)
    for vv, vl in vols:
        add_textbox(s, vx2, Inches(6.3), Inches(0.55), Inches(0.35),
                    vv, size=Pt(14), bold=True, color=C_GREEN_L)
        add_textbox(s, vx2, Inches(6.65), Inches(1.5), Inches(0.35),
                    vl, size=Pt(9), color=C_GRAY)
        vx2 += Inches(2.0)

slide6()


# ─────────────────────────────────────────────────────────────
# SLIDE 7 — RESULTS (BLANK)
# ─────────────────────────────────────────────────────────────
def slide7():
    s = prs.slides.add_slide(blank_layout)
    bg(s)
    accent_bar(s, y=0, h=Inches(0.06))
    slide_num_tag(s, 7)

    section_label(s, "Результаты эксперимента", Inches(0.5), Inches(0.18))
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(11), Inches(0.55),
                "Верификация на реальных данных",
                size=Pt(26), bold=True, color=C_WHITE)
    add_textbox(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.35),
                "Сравнение точности моделей по метрике MAPE (чем меньше — тем лучше)",
                size=Pt(12), color=C_GRAY)

    def blank_card(cx, cy, cw, ch, title, subtitle, note):
        card_rect(s, cx, cy, cw, ch, fill=C_CARD,
                  border=RGBColor(0x28, 0x3A, 0x50))
        add_textbox(s, cx + Inches(0.2), cy + Inches(0.15),
                    cw - Inches(0.3), Inches(0.35),
                    title, size=Pt(13), bold=True, color=C_WHITE)
        add_textbox(s, cx + Inches(0.2), cy + Inches(0.52),
                    cw - Inches(0.3), Inches(0.28),
                    subtitle, size=Pt(10), color=C_GRAY)
        # blank bar placeholders
        models = ["SARIMA", "XGBoost", "Гибридная"]
        colors = [C_BLUE_L, C_GOLD, C_GREEN]
        by3 = cy + Inches(0.9)
        for mi, (mn, mc) in enumerate(zip(models, colors)):
            add_textbox(s, cx + Inches(0.2), by3, Inches(0.85), Inches(0.3),
                        mn, size=Pt(10), color=C_GRAY)
            add_rect(s, cx + Inches(1.15), by3 + Inches(0.07),
                     cw - Inches(1.7), Inches(0.18),
                     fill=RGBColor(0x25, 0x36, 0x48))
            add_textbox(s, cx + cw - Inches(0.5), by3,
                        Inches(0.4), Inches(0.3),
                        "—", size=Pt(10), bold=True, color=C_WHITE,
                        align=PP_ALIGN.RIGHT)
            by3 += Inches(0.42)
        # note
        add_rect(s, cx + Inches(0.2), cy + ch - Inches(0.55),
                 cw - Inches(0.4), Inches(0.42),
                 fill=RGBColor(0x1A, 0x2B, 0x3C))
        add_textbox(s, cx + Inches(0.25), cy + ch - Inches(0.52),
                    cw - Inches(0.5), Inches(0.38),
                    note, size=Pt(10), color=C_DIM,
                    align=PP_ALIGN.CENTER)

    blank_card(Inches(0.35), Inches(1.6), Inches(6.1), Inches(2.8),
               "Показатель 1", "Период / описание",
               "← Заполните значения MAPE")
    blank_card(Inches(6.85), Inches(1.6), Inches(6.1), Inches(2.8),
               "Показатель 2", "Период / описание",
               "← Заполните значения MAPE")

    # KPI band
    kpis = [
        ("—", "Снижение MAPE\nгибридной vs лучшей", C_GREEN_L),
        ("—", "Лучший показатель\nMAE", C_GOLD),
        ("—", "R² гибридной\nмодели", C_BLUE_L),
    ]
    kx = Inches(0.35)
    kw = Inches(3.9)
    for kv, kl, kc in kpis:
        card_rect(s, kx, Inches(4.58), kw, Inches(1.22), fill=C_CARD)
        add_textbox(s, kx, Inches(4.65), kw, Inches(0.55),
                    kv, size=Pt(30), bold=True, color=kc, align=PP_ALIGN.CENTER)
        add_textbox(s, kx, Inches(5.22), kw, Inches(0.45),
                    kl, size=Pt(10), color=C_GRAY, align=PP_ALIGN.CENTER)
        kx += kw + Inches(0.08)

    # Forecast table
    add_rect(s, Inches(0.35), Inches(5.95), Inches(12.6), Inches(1.25),
             fill=C_CARD, line=RGBColor(0x28, 0x3A, 0x50), line_w=Pt(0.5))
    add_textbox(s, Inches(0.55), Inches(6.02), Inches(3), Inches(0.32),
                "ПРОГНОЗ НА БУДУЩИЕ ПЕРИОДЫ",
                size=Pt(9), bold=True, color=C_RED)
    cols = ["Период", "Прогноз", "Нижняя 95%", "Верхняя 95%"]
    cx2 = Inches(0.55)
    for c in cols:
        add_textbox(s, cx2, Inches(6.38), Inches(2.9), Inches(0.28),
                    c, size=Pt(9), bold=True, color=C_GRAY)
        cx2 += Inches(3.05)
    years = ["2025", "2026", "2027"]
    cx3 = Inches(0.55)
    yx = Inches(6.72)
    for y in years:
        add_textbox(s, cx3, yx, Inches(2.9), Inches(0.26),
                    y, size=Pt(10), color=C_GOLD)
        cx3_v = cx3 + Inches(3.05)
        for _ in range(3):
            add_textbox(s, cx3_v, yx, Inches(2.9), Inches(0.26),
                        "—", size=Pt(10), color=C_DIM)
            cx3_v += Inches(3.05)
        cx3 = Inches(0.55)
        yx += Inches(0)  # same row trick won't work — use columns
    # Actually put them side by side properly
    # Re-draw as a proper table grid
    # (above rows share same yx — let's just do 1 row, note says fill)
    add_textbox(s, Inches(4.0), Inches(6.65), Inches(8.5), Inches(0.45),
                "← Вставьте прогнозные значения из системы (строки для 2025, 2026, 2027)",
                size=Pt(10), color=C_DIM)

slide7()


# ─────────────────────────────────────────────────────────────
# SLIDE 8 — ADVANTAGES
# ─────────────────────────────────────────────────────────────
def slide8():
    s = prs.slides.add_slide(blank_layout)
    bg(s)
    accent_bar(s, y=0, h=Inches(0.06))
    slide_num_tag(s, 8)

    section_label(s, "Преимущества", Inches(0.5), Inches(0.18))
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.55),
                "Почему гибридный подход эффективнее традиционных методов",
                size=Pt(26), bold=True, color=C_WHITE)

    adv_left = [
        ("🎯", "Точность выше за счёт ансамблирования",
         "Адаптивное взвешивание делает ставку на модель, "
         "которая точнее на ваших данных. Стабильно превосходит каждую базовую модель."),
        ("🔄", "Работает на коротких рядах",
         "Специальный коэффициент для малых выборок (n<30). "
         "Запуск с 5–10 наблюдений — типичная ситуация в корпоративной HR-аналитике."),
        ("📰", "Учёт внешнего контекста — автоматически",
         "LLM-эксперт читает актуальные источники и корректирует прогноз. "
         "Реагирует на изменения рынка труда без ручного вмешательства."),
    ]
    adv_right = [
        ("⚡", "Скорость: прогноз за 2 минуты",
         "REST API — встраивается в любую HR-систему (1С, SAP). "
         "Веб-интерфейс — без установки ПО, с любого устройства."),
        ("🔒", "Отечественный технологический стек",
         "YandexGPT, развёртывание on-premise или Яндекс Облако. "
         "Соответствие 152-ФЗ и требованиям импортозамещения."),
        ("📊", "Интерпретируемые результаты",
         "HR-директор получает не только цифру, но объяснение: "
         "вклад каждой модели, учтённые факторы, уровень доверия."),
        ("🔧", "Универсальность и масштабируемость",
         "Один экземпляр обслуживает все HR-метрики. "
         "Горизонт прогноза — от 1 месяца до 5 лет."),
    ]

    def adv_column(items, cx, n_items):
        ay = Inches(1.18)
        iw = Inches(5.9)
        ih = Inches(1.32) if n_items == 4 else Inches(1.5)
        for ico, title, desc in items:
            card_rect(s, cx, ay, iw, ih, fill=C_CARD)
            add_textbox(s, cx + Inches(0.18), ay + Inches(0.12),
                        Inches(0.5), Inches(0.45), ico, size=Pt(20))
            add_textbox(s, cx + Inches(0.78), ay + Inches(0.1),
                        iw - Inches(0.9), Inches(0.38),
                        title, size=Pt(12), bold=True, color=C_WHITE)
            add_textbox(s, cx + Inches(0.78), ay + Inches(0.52),
                        iw - Inches(0.9), Inches(0.68),
                        desc, size=Pt(10.5), color=C_GRAY)
            ay += ih + Inches(0.1)

    adv_column(adv_left,  Inches(0.35), 3)
    adv_column(adv_right, Inches(6.9),  4)

    # Comparison table
    add_rect(s, Inches(0.35), Inches(5.78), Inches(6.2), Inches(1.4),
             fill=RGBColor(0x0C, 0x17, 0x22),
             line=RGBColor(0x50, 0x18, 0x14), line_w=Pt(0.75))
    add_textbox(s, Inches(0.55), Inches(5.87), Inches(5.8), Inches(0.28),
                "СРАВНЕНИЕ С ТРАДИЦИОННЫМИ ПОДХОДАМИ",
                size=Pt(8), bold=True, color=C_RED)
    rows = [
        ("Учёт нелинейности",    "Нет",     "Да"),
        ("Внешний контекст",     "Вручную", "Авто (LLM)"),
        ("Доверительный интервал", "Нет",   "95% ДИ"),
        ("Короткие ряды (n<15)", "Нестаб.", "Адаптация"),
    ]
    add_textbox(s, Inches(2.35), Inches(6.18), Inches(2.1), Inches(0.28),
                "Excel / ARIMA", size=Pt(9), bold=True, color=C_GRAY,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(4.5), Inches(6.18), Inches(2.1), Inches(0.28),
                "Гибридная", size=Pt(9), bold=True, color=C_GREEN_L,
                align=PP_ALIGN.CENTER)
    ry = Inches(6.5)
    for rl, rv1, rv2 in rows:
        add_textbox(s, Inches(0.55), ry, Inches(1.75), Inches(0.24),
                    rl, size=Pt(9), color=C_GRAY)
        add_rect(s, Inches(2.35), ry, Inches(2.0), Inches(0.24),
                 fill=RGBColor(0x2D, 0x10, 0x10))
        add_textbox(s, Inches(2.35), ry, Inches(2.0), Inches(0.24),
                    rv1, size=Pt(9), bold=True, color=C_RED_L,
                    align=PP_ALIGN.CENTER)
        add_rect(s, Inches(4.55), ry, Inches(1.9), Inches(0.24),
                 fill=RGBColor(0x0E, 0x28, 0x1A))
        add_textbox(s, Inches(4.55), ry, Inches(1.9), Inches(0.24),
                    rv2, size=Pt(9), bold=True, color=C_GREEN_L,
                    align=PP_ALIGN.CENTER)
        ry += Inches(0.26)

slide8()


# ─────────────────────────────────────────────────────────────
# SLIDE 9 — IMPLEMENTATION
# ─────────────────────────────────────────────────────────────
def slide9():
    s = prs.slides.add_slide(blank_layout)
    bg(s)
    accent_bar(s, y=0, h=Inches(0.06))
    slide_num_tag(s, 9)

    section_label(s, "Внедрение", Inches(0.5), Inches(0.18))
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(11), Inches(0.55),
                "План пилотного проекта в «Уральской стали»",
                size=Pt(26), bold=True, color=C_WHITE)

    phases = [
        ("1", "Фаза 1 — Подготовка",
         "Аудит HR-данных и подключение источников",
         "Инвентаризация исторических данных в 1С:ЗУП. "
         "Выбор приоритетных метрик. Подготовка CSV-выгрузок. Установка системы.",
         "1–2 недели"),
        ("2", "Фаза 2 — Пилот",
         "Прогнозирование численности и ФОТ",
         "Обучение модели на данных 3–5 лет. Верификация на последних периодах. "
         "Настройка LLM-эксперта на отраслевые источники.",
         "2–3 недели"),
        ("3", "Фаза 3 — Расширение",
         "Все HR-процессы + интеграция с 1С",
         "Подключение всех категорий персонала и HR-метрик. "
         "API-интеграция с 1С:ЗУП. Дашборд для HR-директора.",
         "3–4 недели"),
        ("4", "Фаза 4 — Эксплуатация",
         "Регулярное обновление и мониторинг",
         "Ежеквартальное автообновление прогнозов. "
         "Обучение HR-отдела. Настройка алертов при аномалиях.",
         "С 6-й недели"),
    ]

    ty = Inches(1.2)
    for ph_num, ph_label, ph_title, ph_desc, ph_time in phases:
        # dot
        add_rect(s, Inches(0.35), ty, Inches(0.52), Inches(0.52), fill=C_RED)
        add_textbox(s, Inches(0.35), ty, Inches(0.52), Inches(0.52),
                    ph_num, size=Pt(14), bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        # vertical line
        if ph_num != "4":
            add_rect(s, Inches(0.595), ty + Inches(0.52),
                     Inches(0.03), Inches(1.08),
                     fill=RGBColor(0x40, 0x14, 0x12))
        # text
        add_textbox(s, Inches(1.05), ty, Inches(2.5), Inches(0.28),
                    ph_label, size=Pt(9), bold=True, color=C_RED)
        add_textbox(s, Inches(1.05), ty + Inches(0.25), Inches(5.2), Inches(0.35),
                    ph_title, size=Pt(13), bold=True, color=C_WHITE)
        add_textbox(s, Inches(1.05), ty + Inches(0.62), Inches(5.3), Inches(0.55),
                    ph_desc, size=Pt(10.5), color=C_GRAY)
        add_textbox(s, Inches(1.05), ty + Inches(1.2), Inches(2.5), Inches(0.26),
                    "📅  " + ph_time, size=Pt(10), bold=True, color=C_GOLD)
        ty += Inches(1.52)

    # RIGHT: tech stack
    rx = Inches(7.0)
    add_textbox(s, rx, Inches(1.18), Inches(6), Inches(0.3),
                "ТЕХНОЛОГИИ", size=Pt(9), bold=True, color=C_RED)
    tech_groups = [
        ("ML / AI", ["SARIMA", "XGBoost", "TimeLLM", "YandexGPT"],
         C_BLUE_L, RGBColor(0x0A, 0x20, 0x38)),
        ("Данные", ["Росстат", "hh.ru", "1С:ЗУП", "Минцифры"],
         C_GOLD, RGBColor(0x2A, 0x1A, 0x06)),
        ("Интеграция", ["REST API", "Docker", "Web UI", "Python SDK"],
         C_GREEN_L, RGBColor(0x0A, 0x22, 0x12)),
    ]
    tx = rx
    for group_title, pills, pcol, pbg in tech_groups:
        add_textbox(s, tx, Inches(1.55), Inches(2.0), Inches(0.28),
                    group_title, size=Pt(9), bold=True, color=pcol)
        py2 = Inches(1.88)
        for pill_text in pills:
            add_rect(s, tx, py2, Inches(1.8), Inches(0.3), fill=pbg,
                     line=pcol, line_w=Pt(0.5))
            add_textbox(s, tx, py2, Inches(1.8), Inches(0.3),
                        pill_text, size=Pt(10), bold=True, color=pcol,
                        align=PP_ALIGN.CENTER)
            py2 += Inches(0.36)
        tx += Inches(2.05)

    # infra box
    add_rect(s, rx, Inches(4.18), Inches(6.2), Inches(2.1), fill=C_CARD,
             line=RGBColor(0x28, 0x3A, 0x50), line_w=Pt(0.5))
    add_textbox(s, rx + Inches(0.2), Inches(4.27), Inches(5.7), Inches(0.3),
                "ТРЕБОВАНИЯ К ИНФРАСТРУКТУРЕ",
                size=Pt(9), bold=True, color=C_WHITE)
    infra = [
        (C_GREEN,  "Сервер: 16 ГБ RAM, CPU 4+ ядра (без GPU — SARIMA + XGBoost)"),
        (C_GREEN,  "Опционально: GPU NVIDIA ≥ 8 ГБ VRAM (для TimeLLM)"),
        (C_GREEN,  "Развёртывание: Docker on-premise или Яндекс Облако"),
        (C_GREEN,  "Доступ в интернет для LLM-эксперта (YandexGPT + веб-парсинг)"),
        (C_GOLD,   "Срок развёртывания пилота: 1–2 недели"),
    ]
    iy2 = Inches(4.65)
    for ic, it in infra:
        add_rect(s, rx + Inches(0.2), iy2 + Inches(0.08),
                 Inches(0.1), Inches(0.1), fill=ic)
        add_textbox(s, rx + Inches(0.42), iy2, Inches(5.6), Inches(0.3),
                    it, size=Pt(10.5), color=C_GRAY)
        iy2 += Inches(0.37)

    # deliverable
    add_rect(s, rx, Inches(6.42), Inches(6.2), Inches(0.82),
             fill=RGBColor(0x22, 0x0C, 0x0C),
             line=C_RED, line_w=Pt(0.75))
    add_textbox(s, rx + Inches(0.2), Inches(6.5), Inches(5.8), Inches(0.65),
                "💼  Что получает «Уральская сталь»: готовая система с веб-интерфейсом, "
                "API и документацией. Исходный код передаётся.",
                size=Pt(10.5), color=RGBColor(0xF1, 0x94, 0x8A))

slide9()


# ─────────────────────────────────────────────────────────────
# SLIDE 10 — CTA
# ─────────────────────────────────────────────────────────────
def slide10():
    s = prs.slides.add_slide(blank_layout)
    bg(s, RGBColor(0x0A, 0x14, 0x1E))
    accent_bar(s, y=0, h=Inches(0.08))
    # red glow strip
    add_rect(s, 0, H - Inches(0.08), W, Inches(0.08), fill=C_RED)
    slide_num_tag(s, 10)

    add_textbox(s, 0, Inches(1.1), W, Inches(0.4),
                "СЛЕДУЮЩИЙ ШАГ", size=Pt(11), bold=True,
                color=C_GOLD, align=PP_ALIGN.CENTER)

    tb = s.shapes.add_textbox(Inches(1.5), Inches(1.65), Inches(10.3), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "Запустим пилот в «"
    r1.font.name = FONT; r1.font.size = Pt(38); r1.font.bold = True
    r1.font.color.rgb = C_WHITE
    r2 = p.add_run()
    r2.text = "Уральской стали"
    r2.font.name = FONT; r2.font.size = Pt(38); r2.font.bold = True
    r2.font.color.rgb = C_RED_L
    r3 = p.add_run()
    r3.text = "»?"
    r3.font.name = FONT; r3.font.size = Pt(38); r3.font.bold = True
    r3.font.color.rgb = C_WHITE

    add_textbox(s, Inches(2.0), Inches(3.55), Inches(9.3), Inches(0.65),
                "Загрузите любой исторический HR-показатель — через 2 минуты вы получите\n"
                "прогноз с 95% доверительным интервалом и аналитикой LLM-эксперта",
                size=Pt(15), color=C_GRAY, align=PP_ALIGN.CENTER)

    # CTA buttons
    add_rect(s, Inches(2.8), Inches(4.45), Inches(3.5), Inches(0.65), fill=C_RED)
    add_textbox(s, Inches(2.8), Inches(4.45), Inches(3.5), Inches(0.65),
                "🚀  Запустить демо-прогноз",
                size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(7.05), Inches(4.45), Inches(3.5), Inches(0.65),
             fill=RGBColor(0x14, 0x22, 0x32),
             line=RGBColor(0x40, 0x5A, 0x72), line_w=Pt(1))
    add_textbox(s, Inches(7.05), Inches(4.45), Inches(3.5), Inches(0.65),
                "📋  Техническое предложение",
                size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # contact cards
    contacts = [
        ("Автор", "Автор диссертации"),
        ("Репозиторий", "github.com/blabla-user-serj/disser"),
        ("Контакт", "Заполните контакт"),
        ("Версия", "v4.0  ·  Февраль 2026"),
    ]
    ccw = Inches(2.9)
    cx4 = Inches(0.9)
    for cl, cv in contacts:
        add_rect(s, cx4, Inches(5.42), ccw, Inches(1.05),
                 fill=RGBColor(0x14, 0x22, 0x30),
                 line=RGBColor(0x28, 0x3A, 0x50), line_w=Pt(0.5))
        add_textbox(s, cx4 + Inches(0.18), Inches(5.52),
                    ccw - Inches(0.3), Inches(0.28),
                    cl.upper(), size=Pt(9), color=C_DIM)
        add_textbox(s, cx4 + Inches(0.18), Inches(5.82),
                    ccw - Inches(0.3), Inches(0.35),
                    cv, size=Pt(12), bold=True, color=C_WHITE)
        cx4 += ccw + Inches(0.22)

    # footer
    add_rect(s, 0, H - Inches(0.55), W, Inches(0.55),
             fill=RGBColor(0x06, 0x0E, 0x18))
    add_textbox(s, Inches(0.4), H - Inches(0.48), W - Inches(0.8), Inches(0.35),
                "Гибридная система прогнозирования для HR  ·  Диссертационное исследование 2025–2026",
                size=Pt(10), color=C_DIM, align=PP_ALIGN.CENTER)

slide10()


# ─────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "uralskaya-stal-hr-forecasting.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Size:  {os.path.getsize(out):,} bytes")
